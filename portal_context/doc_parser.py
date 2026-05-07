"""
Document Parser — Parses supplementary documents into plain text.

Supports PDF, DOCX, PPTX, MD, and TXT files.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def parse_documents(file_paths: list[str]) -> str:
    """
    Parse multiple documents and return combined text.
    
    Args:
        file_paths: List of file paths to parse
        
    Returns:
        Combined text content from all documents
    """
    all_text = []
    for path in file_paths:
        try:
            text = parse_single_document(path)
            if text.strip():
                all_text.append(f"--- Document: {Path(path).name} ---\n{text}")
                logger.info(f"Parsed document: {path} ({len(text)} chars)")
        except Exception as e:
            logger.warning(f"Failed to parse {path}: {e}")
    return "\n\n".join(all_text)


def parse_single_document(file_path: str) -> str:
    """Parse a single document file to text."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".md": _parse_text,
        ".txt": _parse_text,
        ".markdown": _parse_text,
        ".rst": _parse_text,
    }

    parser = parsers.get(suffix)
    if parser is None:
        logger.warning(f"Unsupported file type: {suffix}. Trying as plain text.")
        return _parse_text(file_path)

    return parser(file_path)


def _parse_pdf(file_path: str) -> str:
    """Parse PDF file to text."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        raise ImportError("PyPDF2 is required for PDF parsing: pip install PyPDF2")

    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def _parse_docx(file_path: str) -> str:
    """Parse DOCX file to text."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required for DOCX parsing: pip install python-docx")

    doc = Document(file_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n\n".join(paragraphs)


def _parse_pptx(file_path: str) -> str:
    """Parse PPTX file to text."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PPTX parsing: pip install python-pptx")

    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


def _parse_text(file_path: str) -> str:
    """Parse plain text / markdown file."""
    return Path(file_path).read_text(encoding="utf-8", errors="replace")
